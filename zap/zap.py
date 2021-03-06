import requests
import argparse
import wikipedia as wp
from pptx import Presentation
from pptx.util import Inches, Pt


class zap:

    """
    Zap class for providing interface with the software
    """

    def __init__(self):
        """
        Initialize for the zap class, creates a parser object and uses it to
        parse options
        """
        self.desc = "Python based Wikipedia article's PPT generator"
        self.parser = argparse.ArgumentParser(description=(self.desc))
        self.parser.add_argument(
            '-t', dest='term', type=str, default=None, help='Term to search')
        args = vars(self.parser.parse_args())

        if args['term'] is None:
            print 'Passing a term is necessary'
            self.parser.print_help()
            exit(-1)
        self.searchForTerm(args['term'])

    def searchForTerm(self, term):
        """
        Searches for term on wikipedia to find the related pages
        term: The argument provided by the user to be searched
        searchResults: The returned list of search results
        """
        print '[+] Searching for ' + term + ' at Wikipedia'
        searchResults = wp.search(term)
        print '[+] We found following results:'
        for i in range(0, len(searchResults)):
            print str(i) + ') ' + searchResults[i]
        selected = raw_input(
            'Enter the no for which you want to generate Slides:')
        self.getContent(searchResults[int(selected)])

    def getContent(self, pageName):
        """
        Gets the content from a page
        pageName: The page title to be retrieved
        page: The returned object with title, content and links
        """
        page = wp.page(pageName)
        self.generatePPT(page)

    def generatePPT(self, page):
        """
        This function generates the PPT by adding various kind of slides
        throughout its execution
        """
        self.prs = Presentation()
        layout = self.prs.slide_layouts[6]

    def addSlide(self, layout, content):
        """
        To add a single kind of slide
        """
        slide = self.prs.slides.add_slide(layout)

        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.textframe
        tf.text = content


def main():
    zap()

if __name__ == '__main__':
    main()
